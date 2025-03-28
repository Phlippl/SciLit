"""
Dokumentenverarbeitung für SciLit
---------------------------------
Dieses Modul enthält alle Funktionen zur Verarbeitung von Dokumenten,
einschließlich Extraktion von Text und Metadaten aus verschiedenen Dateiformaten
und Abrufen von zusätzlichen Metadaten über verschiedene APIs.
"""

import os
import logging
import tempfile
from typing import Dict, List, Optional, Tuple, Any
import re
from pathlib import Path
import shutil
import json
import hashlib

# Dokumentenverarbeitungs-Bibliotheken
import PyPDF2
import pytesseract
from PIL import Image
import docx
import mammoth
import spacy
import epub2txt

# Metadaten-APIs
import requests
from langchain.text_splitter import RecursiveCharacterTextSplitter

# Logger konfigurieren
logger = logging.getLogger("scilit.document_processor")

# Konfigurieren der Modelle und APIs
SPACY_MODEL_DE = "de_core_news_sm"
SPACY_MODEL_EN = "en_core_web_sm"

# OCR Konfiguration
OCR_LANGUAGES = {
    'de': 'deu',
    'en': 'eng',
    'auto': 'deu+eng'  # Mehrsprachenerkennung
}

class DocumentProcessor:
    """Hauptklasse für die Verarbeitung von Dokumenten."""
    
    def __init__(
        self, 
        upload_dir: str, 
        processed_dir: str,
        ocr_if_needed: bool = True,
        language: str = "auto",
        metadata_sources: Dict[str, bool] = None
    ):
        """
        Initialisiert den Dokumentenprozessor.
        
        Args:
            upload_dir: Verzeichnis für hochgeladene Dateien
            processed_dir: Verzeichnis für verarbeitete Dateien
            ocr_if_needed: Ob OCR bei Bedarf angewendet werden soll
            language: Hauptsprache der Dokumente ('de', 'en', 'auto', 'mixed')
            metadata_sources: Welche Metadaten-APIs verwendet werden sollen
        """
        self.upload_dir = upload_dir
        self.processed_dir = processed_dir
        self.ocr_if_needed = ocr_if_needed
        self.language = language
        
        # Standardwerte für Metadatenquellen, wenn nicht angegeben
        self.metadata_sources = metadata_sources or {
            "use_crossref": True,
            "use_openlib": True,
            "use_k10plus": True,
            "use_googlebooks": True,
            "use_openalex": True,
        }
        
        # Verzeichnisse sicherstellen
        os.makedirs(self.upload_dir, exist_ok=True)
        os.makedirs(self.processed_dir, exist_ok=True)
        
        # SpaCy Modelle laden (verzögert, um nicht unnötig Speicher zu belegen)
        self._nlp_de = None
        self._nlp_en = None
        
        # Textsplitter für Chunks erstellen
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
        )
    
    def _load_spacy_model(self, language: str):
        """Lädt das entsprechende SpaCy-Modell nach Bedarf."""
        if language == "de" and self._nlp_de is None:
            try:
                self._nlp_de = spacy.load(SPACY_MODEL_DE)
                logger.info(f"SpaCy Modell '{SPACY_MODEL_DE}' geladen")
            except OSError:
                logger.warning(f"SpaCy Modell '{SPACY_MODEL_DE}' nicht gefunden, wird heruntergeladen...")
                spacy.cli.download(SPACY_MODEL_DE)
                self._nlp_de = spacy.load(SPACY_MODEL_DE)
        elif language == "en" and self._nlp_en is None:
            try:
                self._nlp_en = spacy.load(SPACY_MODEL_EN)
                logger.info(f"SpaCy Modell '{SPACY_MODEL_EN}' geladen")
            except OSError:
                logger.warning(f"SpaCy Modell '{SPACY_MODEL_EN}' nicht gefunden, wird heruntergeladen...")
                spacy.cli.download(SPACY_MODEL_EN)
                self._nlp_en = spacy.load(SPACY_MODEL_EN)
        
        return self._nlp_de if language == "de" else self._nlp_en
    
    def process_document(self, filename: str) -> Dict[str, Any]:
        """
        Verarbeitet ein Dokument vollständig.
        
        Args:
            filename: Name der zu verarbeitenden Datei im Upload-Verzeichnis
            
        Returns:
            Dict mit Metadaten und Verarbeitungsergebnissen
        """
        filepath = os.path.join(self.upload_dir, filename)
        logger.info(f"Verarbeite Dokument: {filepath}")
        
        # Eindeutige ID für das Dokument erzeugen
        doc_id = self._generate_document_id(filepath)
        
        # Metadaten-Verzeichnis erstellen
        doc_dir = os.path.join(self.processed_dir, doc_id)
        os.makedirs(doc_dir, exist_ok=True)
        
        # Datei ins verarbeitete Verzeichnis kopieren
        target_filepath = os.path.join(doc_dir, filename)
        shutil.copy2(filepath, target_filepath)
        
        # Extrahiere Text und Basis-Metadaten aus dem Dokument
        text, basic_metadata = self.extract_content_and_metadata(filepath)
        
        # Erweiterte Metadaten über APIs abrufen
        metadata = self.enhance_metadata(basic_metadata)
        
        # Text in Chunks aufteilen
        chunks = self.split_text_into_chunks(text, metadata["language"] if "language" in metadata else self.language)
        
        # Ergebnisse speichern
        self._save_processing_results(doc_dir, text, metadata, chunks)
        
        # Rückgabe-Dictionary erstellen
        result = {
            "id": doc_id,
            "filename": filename,
            "filepath": target_filepath,
            "metadata": metadata,
            "chunks_count": len(chunks),
            "status": "processed"
        }
        
        logger.info(f"Dokument erfolgreich verarbeitet: {doc_id}")
        return result
    
    def _generate_document_id(self, filepath: str) -> str:
        """Erzeugt eine eindeutige ID für ein Dokument basierend auf Inhalt und Name."""
        filename = os.path.basename(filepath)
        file_stats = os.stat(filepath)
        
        # Kombiniere Dateiname, Größe und Änderungszeit für einen eindeutigen Hashwert
        id_string = f"{filename}_{file_stats.st_size}_{file_stats.st_mtime}"
        return hashlib.md5(id_string.encode()).hexdigest()[:12]
    
    def extract_content_and_metadata(self, filepath: str) -> Tuple[str, Dict[str, Any]]:
        """
        Extrahiert Text und grundlegende Metadaten aus einem Dokument.
        
        Args:
            filepath: Pfad zur Datei
            
        Returns:
            Tuple aus (extrahierter Text, Dictionary mit Metadaten)
        """
        file_ext = os.path.splitext(filepath)[1].lower()
        
        # Unterschiedliche Verarbeitung je nach Dateityp
        if file_ext == '.pdf':
            return self._process_pdf(filepath)
        elif file_ext == '.docx':
            return self._process_docx(filepath)
        elif file_ext == '.txt':
            return self._process_txt(filepath)
        elif file_ext == '.epub':
            return self._process_epub(filepath)
        elif file_ext in ['.pptx', '.ppt']:
            return self._process_pptx(filepath)
        else:
            logger.warning(f"Dateityp {file_ext} wird nicht unterstützt")
            raise ValueError(f"Dateityp {file_ext} wird nicht unterstützt")
    
    def _process_pdf(self, filepath: str) -> Tuple[str, Dict[str, Any]]:
        """Verarbeitet eine PDF-Datei."""
        text = ""
        metadata = {}
        
        try:
            with open(filepath, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                # Metadaten aus PDF extrahieren
                pdf_info = pdf_reader.metadata
                if pdf_info:
                    # Standardmetadaten extrahieren
                    for key in pdf_info:
                        clean_key = key.strip('/').lower()
                        metadata[clean_key] = pdf_info[key]
                
                # Text aus allen Seiten extrahieren
                page_texts = []
                needs_ocr = False
                
                for i, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    
                    # Überprüfen, ob OCR benötigt wird (wenn Seite nicht genug extrahierbaren Text hat)
                    if not page_text or len(page_text.strip()) < 100:
                        if self.ocr_if_needed:
                            logger.info(f"Seite {i+1} hat wenig Text - versuche OCR")
                            needs_ocr = True
                            break
                    
                    page_texts.append(page_text)
                
                # OCR anwenden, wenn nötig und aktiviert
                if needs_ocr and self.ocr_if_needed:
                    logger.info(f"Wende OCR auf {filepath} an")
                    text = self._apply_ocr_to_pdf(filepath)
                else:
                    text = "\n\n".join(page_texts)
                
                # Extrahiere Titel und Autor aus dem Text, wenn nicht in Metadaten
                if 'title' not in metadata or not metadata['title']:
                    potential_title = self._extract_title_from_text(text)
                    if potential_title:
                        metadata['title'] = potential_title
                
                if 'author' not in metadata or not metadata['author']:
                    potential_authors = self._extract_authors_from_text(text)
                    if potential_authors:
                        metadata['author'] = potential_authors
                
                # Versuche, das Publikationsjahr zu extrahieren
                if 'creationdate' in metadata:
                    year_match = re.search(r'D:(\d{4})', metadata['creationdate'])
                    if year_match:
                        metadata['year'] = int(year_match.group(1))
                
                # Seitenzahl hinzufügen
                metadata['page_count'] = len(pdf_reader.pages)
                
        except Exception as e:
            logger.error(f"Fehler beim Verarbeiten von PDF {filepath}: {str(e)}")
            raise
        
        return text, metadata
    
    def _apply_ocr_to_pdf(self, filepath: str) -> str:
        """Wendet OCR auf eine PDF-Datei an."""
        
        logger.info(f"Starte OCR-Verarbeitung für {filepath}")
        
        # Versuche zuerst mit PyMuPDF (fitz)
        try:
            import fitz  # PyMuPDF
            
            ocr_lang = OCR_LANGUAGES.get(self.language, OCR_LANGUAGES['auto'])
            doc = fitz.open(filepath)
            text_pages = []
            
            with tempfile.TemporaryDirectory() as temp_dir:
                for page_num, page in enumerate(doc):
                    logger.info(f"Verarbeite Seite {page_num+1}/{len(doc)}")
                    # Bild der Seite extrahieren
                    pix = page.get_pixmap()
                    img_path = f"{temp_dir}/page_{page_num+1}.png"
                    pix.save(img_path)
                    
                    # OCR anwenden
                    try:
                        img = Image.open(img_path)
                        page_text = pytesseract.image_to_string(img, lang=ocr_lang)
                        text_pages.append(page_text)
                    except Exception as e:
                        logger.error(f"OCR-Fehler bei Seite {page_num+1}: {str(e)}")
                        text_pages.append(f"[OCR-FEHLER: {str(e)}]")
            
            return "\n\n".join(text_pages)
            
        except ImportError as e:
            logger.warning(f"PyMuPDF (fitz) nicht verfügbar: {str(e)}. Versuche alternative Methode...")
            
        # Alternative Methode mit pdf2image
        try:
            from pdf2image import convert_from_path
            
            ocr_lang = OCR_LANGUAGES.get(self.language, OCR_LANGUAGES['auto'])
            text_pages = []
            
            with tempfile.TemporaryDirectory() as temp_dir:
                try:
                    # PDF in Bilder umwandeln
                    logger.info(f"Konvertiere PDF zu Bildern mit pdf2image...")
                    pages = convert_from_path(filepath, dpi=300, output_folder=temp_dir)
                    
                    for i, page in enumerate(pages):
                        logger.info(f"OCR für Seite {i+1}/{len(pages)}")
                        # OCR für jede Seite anwenden
                        try:
                            page_text = pytesseract.image_to_string(page, lang=ocr_lang)
                            text_pages.append(page_text)
                        except Exception as e:
                            logger.error(f"OCR-Fehler bei Seite {i+1}: {str(e)}")
                            text_pages.append(f"[OCR-FEHLER: {str(e)}]")
                    
                    return "\n\n".join(text_pages)
                except Exception as e:
                    logger.error(f"Fehler bei pdf2image: {str(e)}")
                    # Fallback-Text zurückgeben mit Fehlermeldung
                    return f"[FEHLER BEI DER OCR-VERARBEITUNG: {str(e)}. Bitte stelle sicher, dass Poppler korrekt installiert ist.]"
                    
        except ImportError as e:
            error_msg = f"Weder PyMuPDF (fitz) noch pdf2image sind verfügbar: {str(e)}"
            logger.error(error_msg)
            return f"[OCR NICHT VERFÜGBAR: {error_msg}. Bitte installiere PyMuPDF oder pdf2image.]"
    
    def _process_docx(self, filepath: str) -> Tuple[str, Dict[str, Any]]:
        """Verarbeitet eine DOCX-Datei."""
        doc = docx.Document(filepath)
        text = "\n".join([para.text for para in doc.paragraphs])
        
        # Metadaten extrahieren
        metadata = {}
        core_properties = doc.core_properties
        
        if core_properties.title:
            metadata['title'] = core_properties.title
        if core_properties.author:
            metadata['author'] = core_properties.author
        if core_properties.created:
            metadata['created'] = core_properties.created.isoformat()
            metadata['year'] = core_properties.created.year
        
        # Seitenzahl ist in DOCX nicht direkt verfügbar, kann geschätzt werden
        # Alternativ mit mammoth HTML erzeugen und Seitenzahl schätzen
        
        # Wenn keine Metadaten gefunden wurden, Text nach Titel/Autor durchsuchen
        if 'title' not in metadata or not metadata['title']:
            potential_title = self._extract_title_from_text(text)
            if potential_title:
                metadata['title'] = potential_title
        
        if 'author' not in metadata or not metadata['author']:
            potential_authors = self._extract_authors_from_text(text)
            if potential_authors:
                metadata['author'] = potential_authors
        
        return text, metadata
    
    def _process_txt(self, filepath: str) -> Tuple[str, Dict[str, Any]]:
        """Verarbeitet eine TXT-Datei."""
        with open(filepath, 'r', encoding='utf-8', errors='replace') as file:
            text = file.read()
        
        # Bei TXT-Dateien müssen wir Metadaten aus dem Text extrahieren
        metadata = {}
        
        # Erste nicht-leere Zeile könnte der Titel sein
        lines = [line.strip() for line in text.split('\n') if line.strip()]
        if lines:
            metadata['title'] = lines[0]
        
        # Versuche, Autoren zu extrahieren
        potential_authors = self._extract_authors_from_text(text)
        if potential_authors:
            metadata['author'] = potential_authors
        
        return text, metadata
    
    def _process_epub(self, filepath: str) -> Tuple[str, Dict[str, Any]]:
        """Verarbeitet eine EPUB-Datei."""
        # epub2txt extrahiert Text aus EPUB-Dateien
        text = epub2txt.epub2txt(filepath)
        
        # Metadaten-Extraktion aus EPUB erfordert zusätzliche Bibliotheken
        metadata = {}
        
        try:
            import ebooklib
            from ebooklib import epub
            
            book = epub.read_epub(filepath)
            
            # Metadaten extrahieren
            if book.metadata:
                # Dublin Core Metadaten
                dc_metadata = book.metadata.get('http://purl.org/dc/elements/1.1/', {})
                
                # Titel
                if 'title' in dc_metadata:
                    metadata['title'] = dc_metadata['title'][0][0]
                
                # Autoren
                if 'creator' in dc_metadata:
                    authors = [creator[0] for creator in dc_metadata['creator']]
                    metadata['author'] = authors
                
                # Publikationsdatum
                if 'date' in dc_metadata:
                    date_str = dc_metadata['date'][0][0]
                    # Versuche, Jahr zu extrahieren
                    year_match = re.search(r'(\d{4})', date_str)
                    if year_match:
                        metadata['year'] = int(year_match.group(1))
                
                # Publisher
                if 'publisher' in dc_metadata:
                    metadata['publisher'] = dc_metadata['publisher'][0][0]
        
        except ImportError:
            logger.warning("ebooklib nicht installiert, beschränkte Metadaten-Extraktion")
            
            # Fallback: Extrahiere Metadaten aus dem Text
            potential_title = self._extract_title_from_text(text)
            if potential_title:
                metadata['title'] = potential_title
            
            potential_authors = self._extract_authors_from_text(text)
            if potential_authors:
                metadata['author'] = potential_authors
        
        return text, metadata
    
    def _process_pptx(self, filepath: str) -> Tuple[str, Dict[str, Any]]:
        """Verarbeitet eine PPTX-Datei."""
        try:
            from pptx import Presentation
            
            presentation = Presentation(filepath)
            
            # Text aus allen Folien extrahieren
            text_runs = []
            
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            
            text = "\n\n".join(text_runs)
            
            # Metadaten extrahieren
            metadata = {}
            core_properties = presentation.core_properties
            
            if core_properties.title:
                metadata['title'] = core_properties.title
            if core_properties.author:
                metadata['author'] = core_properties.author
            if core_properties.created:
                metadata['created'] = core_properties.created.isoformat()
                metadata['year'] = core_properties.created.year
            
            # Anzahl der Folien
            metadata['slide_count'] = len(presentation.slides)
            
            # Wenn keine Metadaten gefunden wurden, Text nach Titel/Autor durchsuchen
            if 'title' not in metadata or not metadata['title']:
                potential_title = self._extract_title_from_text(text)
                if potential_title:
                    metadata['title'] = potential_title
            
            if 'author' not in metadata or not metadata['author']:
                potential_authors = self._extract_authors_from_text(text)
                if potential_authors:
                    metadata['author'] = potential_authors
            
            return text, metadata
            
        except ImportError:
            logger.error("python-pptx nicht installiert")
            raise ImportError("Zur Verarbeitung von PPTX-Dateien wird python-pptx benötigt")
    
    def _extract_title_from_text(self, text: str) -> Optional[str]:
        """Extrahiert einen möglichen Titel aus dem Text."""
        # Erste paar Zeilen des Texts durchsuchen
        lines = [line.strip() for line in text.split('\n') if line.strip()][:10]
        
        for line in lines:
            # Titel sind typischerweise kurz, aber nicht zu kurz, und enden nicht mit Punkt
            if 3 < len(line) < 200 and not line.endswith('.'):
                return line
        
        return None
    
    def _extract_authors_from_text(self, text: str) -> Optional[List[str]]:
        """Extrahiert mögliche Autoren aus dem Text."""
        # Häufige Muster für Autorenlisten
        patterns = [
            r'(?:Autor(?:en)?|Author(?:s)?|von|by)[:\s]+([^\n\.]{3,100})',
            r'(?:^|\n)([A-Z][a-z]+(?:\s[A-Z][a-z]+)+(?:,\s*[A-Z][a-z]+(?:\s[A-Z][a-z]+)+)*)'
        ]
        
        for pattern in patterns:
            matches = re.search(pattern, text[:1000])
            if matches:
                authors_text = matches.group(1)
                
                # Trenne Autoren, wenn sie durch Kommas oder "und"/"and" getrennt sind
                authors = re.split(r',\s*|\s+(?:und|and|&)\s+', authors_text)
                
                # Bereinige die Autorenliste
                authors = [author.strip() for author in authors if len(author.strip()) > 3]
                
                if authors:
                    return authors
        
        return None
    
    def enhance_metadata(self, basic_metadata: Dict[str, Any]) -> Dict[str, Any]:
        """
        Erweitert die grundlegenden Metadaten durch Abfragen verschiedener APIs.
        
        Args:
            basic_metadata: Aus dem Dokument extrahierte Metadaten
            
        Returns:
            Erweiterte Metadaten
        """
        metadata = basic_metadata.copy()
        
        # Titel oder Autor aus den Basis-Metadaten extrahieren
        title = metadata.get('title', '')
        authors = metadata.get('author', [])
        if isinstance(authors, str):
            authors = [authors]
        
        # DOI extrahieren, falls vorhanden
        doi = None
        for key, value in metadata.items():
            if key.lower() in ['doi', 'identifier']:
                if isinstance(value, str):
                    doi_match = re.search(r'10\.\d{4,9}/[-._;()/:a-zA-Z0-9]+', value)
                    if doi_match:
                        doi = doi_match.group(0)
                        break
        
        # 1. CrossRef API
        if self.metadata_sources.get("use_crossref", True) and (title or authors or doi):
            try:
                crossref_metadata = self._fetch_crossref_metadata(title, authors, doi)
                if crossref_metadata:
                    logger.info("Metadaten über CrossRef gefunden")
                    metadata.update(crossref_metadata)
            except Exception as e:
                logger.warning(f"Fehler bei CrossRef-Abfrage: {str(e)}")
        
        # 2. Open Library API
        if self.metadata_sources.get("use_openlib", True) and (title or authors):
            try:
                openlib_metadata = self._fetch_openlib_metadata(title, authors)
                if openlib_metadata:
                    logger.info("Metadaten über Open Library gefunden")
                    metadata.update(openlib_metadata)
            except Exception as e:
                logger.warning(f"Fehler bei Open Library-Abfrage: {str(e)}")
        
        # 3. K10plus API (falls Implementierung vorhanden)
        if self.metadata_sources.get("use_k10plus", True) and (title or authors):
            try:
                k10plus_metadata = self._fetch_k10plus_metadata(title, authors)
                if k10plus_metadata:
                    logger.info("Metadaten über K10plus gefunden")
                    metadata.update(k10plus_metadata)
            except Exception as e:
                logger.warning(f"Fehler bei K10plus-Abfrage: {str(e)}")
        
        # 4. Google Books API
        if self.metadata_sources.get("use_googlebooks", True) and (title or authors):
            try:
                googlebooks_metadata = self._fetch_googlebooks_metadata(title, authors)
                if googlebooks_metadata:
                    logger.info("Metadaten über Google Books gefunden")
                    metadata.update(googlebooks_metadata)
            except Exception as e:
                logger.warning(f"Fehler bei Google Books-Abfrage: {str(e)}")
        
        # 5. OpenAlex API
        if self.metadata_sources.get("use_openalex", True) and (title or authors):
            try:
                openalex_metadata = self._fetch_openalex_metadata(title, authors)
                if openalex_metadata:
                    logger.info("Metadaten über OpenAlex gefunden")
                    metadata.update(openalex_metadata)
            except Exception as e:
                logger.warning(f"Fehler bei OpenAlex-Abfrage: {str(e)}")
        
        return metadata
    
    def _fetch_crossref_metadata(self, title: str, authors: List[str], doi: Optional[str]) -> Dict[str, Any]:
        """Ruft Metadaten von CrossRef ab."""
        metadata = {}
        
        if doi:
            # Direkte Suche nach DOI
            url = f"https://api.crossref.org/works/{doi}"
            response = requests.get(url)
            
            if response.status_code == 200:
                data = response.json()
                if 'message' in data:
                    message = data['message']
                    
                    # Titel
                    if 'title' in message and message['title']:
                        metadata['title'] = message['title'][0]
                    
                    # Autoren
                    if 'author' in message:
                        authors = []
                        for author in message['author']:
                            name_parts = []
                            if 'given' in author:
                                name_parts.append(author['given'])
                            if 'family' in author:
                                name_parts.append(author['family'])
                            if name_parts:
                                authors.append(' '.join(name_parts))
                        if authors:
                            metadata['author'] = authors
                    
                    # Journal/Quelle
                    if 'container-title' in message and message['container-title']:
                        metadata['journal'] = message['container-title'][0]
                    
                    # Jahr
                    if 'published' in message and 'date-parts' in message['published']:
                        date_parts = message['published']['date-parts'][0]
                        if date_parts and len(date_parts) > 0:
                            metadata['year'] = date_parts[0]
                    
                    # Verleger
                    if 'publisher' in message:
                        metadata['publisher'] = message['publisher']
                    
                    # DOI
                    metadata['doi'] = doi
                    
                    return metadata
        
        # Suche nach Titel und/oder Autor, wenn kein DOI oder DOI-Suche erfolglos
        if title or authors:
            query_parts = []
            
            if title:
                query_parts.append(f"title:{title}")
            
            if authors:
                # Verwende nur den ersten Autor für die Suche
                query_parts.append(f"author:{authors[0]}")
            
            query = " ".join(query_parts)
            url = f"https://api.crossref.org/works?query={query}&rows=1"
            
            response = requests.get(url)
            
            if response.status_code == 200:
                data = response.json()
                if 'message' in data and 'items' in data['message'] and data['message']['items']:
                    item = data['message']['items'][0]
                    
                    # Titel
                    if 'title' in item and item['title']:
                        metadata['title'] = item['title'][0]
                    
                    # Autoren
                    if 'author' in item:
                        authors = []
                        for author in item['author']:
                            name_parts = []
                            if 'given' in author:
                                name_parts.append(author['given'])
                            if 'family' in author:
                                name_parts.append(author['family'])
                            if name_parts:
                                authors.append(' '.join(name_parts))
                        if authors:
                            metadata['author'] = authors
                    
                    # Journal/Quelle
                    if 'container-title' in item and item['container-title']:
                        metadata['journal'] = item['container-title'][0]
                    
                    # Jahr
                    if 'published' in item and 'date-parts' in item['published']:
                        date_parts = item['published']['date-parts'][0]
                        if date_parts and len(date_parts) > 0:
                            metadata['year'] = date_parts[0]
                    
                    # Verleger
                    if 'publisher' in item:
                        metadata['publisher'] = item['publisher']
                    
                    # DOI
                    if 'DOI' in item:
                        metadata['doi'] = item['DOI']
                    
                    return metadata
        
        return metadata
    
    def _fetch_openlib_metadata(self, title: str, authors: List[str]) -> Dict[str, Any]:
        """Ruft Metadaten von Open Library ab."""
        metadata = {}
        
        if not title:
            return metadata
        
        # Suche nach Titel
        query = title
        if authors and authors[0]:
            # Füge ersten Autor zur Suche hinzu
            query += f" author:{authors[0]}"
        
        url = f"https://openlibrary.org/search.json?q={query}"
        response = requests.get(url)
        
        if response.status_code == 200:
            data = response.json()
            if 'docs' in data and data['docs']:
                doc = data['docs'][0]
                
                # Titel
                if 'title' in doc:
                    metadata['title'] = doc['title']
                
                # Autoren
                if 'author_name' in doc:
                    metadata['author'] = doc['author_name']
                
                # Jahr
                if 'first_publish_year' in doc:
                    metadata['year'] = doc['first_publish_year']
                
                # Verleger
                if 'publisher' in doc:
                    metadata['publisher'] = doc['publisher'][0] if isinstance(doc['publisher'], list) and doc['publisher'] else doc['publisher']
                
                # ISBN
                if 'isbn' in doc and doc['isbn']:
                    metadata['isbn'] = doc['isbn'][0] if isinstance(doc['isbn'], list) else doc['isbn']
                
                return metadata
        
        return metadata
    
    def _fetch_k10plus_metadata(self, title: str, authors: List[str]) -> Dict[str, Any]:
        """Ruft Metadaten vom K10plus-Katalog ab."""
        metadata = {}
        
        # K10plus verwendet SRU (Search/Retrieve via URL)
        # Beispiel: https://sru.k10plus.de/opac-de-627?version=1.1&operation=searchRetrieve&query=pica.tit%3D%22Digitalisierung%22
        
        if not title:
            return metadata
        
        # Bereite die Suche vor (nur Titel, da das API keine kombinierte Suche unterstützt)
        query = f'pica.tit="{title}"'
        url = f"https://sru.k10plus.de/opac-de-627?version=1.1&operation=searchRetrieve&query={query}&maximumRecords=1"
        
        response = requests.get(url)
        
        if response.status_code == 200:
            # Das Ergebnis ist in XML - wir benötigen elementtree zur Verarbeitung
            import xml.etree.ElementTree as ET
            
            try:
                root = ET.fromstring(response.text)
                
                # Namespace-Definition
                ns = {
                    'zs': 'http://www.loc.gov/zing/srw/',
                    'pica': 'info:srw/schema/5/picaXML-v1.0',
                    'marc': 'http://www.loc.gov/MARC21/slim'
                }
                
                # Prüfe, ob Ergebnisse vorhanden sind
                record_count = root.find('.//zs:numberOfRecords', ns)
                if record_count is not None and record_count.text != '0':
                    # Finde den ersten Datensatz
                    record = root.find('.//pica:record', ns)
                    
                    # Titel extrahieren
                    title_field = record.find('.//pica:datafield[@tag="021A"]/pica:subfield[@code="a"]', ns)
                    if title_field is not None:
                        metadata['title'] = title_field.text
                    
                    # Autoren extrahieren
                    authors = []
                    author_fields = record.findall('.//pica:datafield[@tag="028A"]', ns)
                    for author_field in author_fields:
                        author_parts = []
                        
                        # Nachname
                        last_name = author_field.find('./pica:subfield[@code="a"]', ns)
                        if last_name is not None:
                            author_parts.append(last_name.text)
                        
                        # Vorname
                        first_name = author_field.find('./pica:subfield[@code="d"]', ns)
                        if first_name is not None:
                            author_parts.append(first_name.text)
                        
                        if author_parts:
                            authors.append(' '.join(author_parts))
                    
                    if authors:
                        metadata['author'] = authors
                    
                    # Jahr extrahieren
                    year_field = record.find('.//pica:datafield[@tag="011@"]/pica:subfield[@code="a"]', ns)
                    if year_field is not None:
                        # Jahr aus Datumsstring extrahieren
                        year_match = re.search(r'(\d{4})', year_field.text)
                        if year_match:
                            metadata['year'] = int(year_match.group(1))
                    
                    # Verlag extrahieren
                    publisher_field = record.find('.//pica:datafield[@tag="033A"]/pica:subfield[@code="n"]', ns)
                    if publisher_field is not None:
                        metadata['publisher'] = publisher_field.text
                    
                    # ISBN extrahieren
                    isbn_field = record.find('.//pica:datafield[@tag="004A"]/pica:subfield[@code="0"]', ns)
                    if isbn_field is not None:
                        metadata['isbn'] = isbn_field.text
                    
                    return metadata
            
            except ET.ParseError:
                logger.warning("Fehler beim Parsen der K10plus XML-Antwort")
        
        return metadata
    
    def _fetch_googlebooks_metadata(self, title: str, authors: List[str]) -> Dict[str, Any]:
        """Ruft Metadaten von Google Books ab."""
        metadata = {}
        
        if not title:
            return metadata
        
        # Bereite die Suche vor
        query = f"intitle:{title}"
        if authors and authors[0]:
            query += f"+inauthor:{authors[0]}"
        
        url = f"https://www.googleapis.com/books/v1/volumes?q={query}&maxResults=1"
        response = requests.get(url)
        
        if response.status_code == 200:
            data = response.json()
            if 'items' in data and data['items']:
                item = data['items'][0]
                volume_info = item.get('volumeInfo', {})
                
                # Titel
                if 'title' in volume_info:
                    metadata['title'] = volume_info['title']
                
                # Autoren
                if 'authors' in volume_info:
                    metadata['author'] = volume_info['authors']
                
                # Verlag
                if 'publisher' in volume_info:
                    metadata['publisher'] = volume_info['publisher']
                
                # Jahr
                if 'publishedDate' in volume_info:
                    # Jahr aus Datumsstring extrahieren
                    year_match = re.search(r'(\d{4})', volume_info['publishedDate'])
                    if year_match:
                        metadata['year'] = int(year_match.group(1))
                
                # ISBN
                if 'industryIdentifiers' in volume_info:
                    for identifier in volume_info['industryIdentifiers']:
                        if identifier['type'] in ['ISBN_13', 'ISBN_10']:
                            metadata['isbn'] = identifier['identifier']
                            break
                
                # Seitenzahl
                if 'pageCount' in volume_info:
                    metadata['page_count'] = volume_info['pageCount']
                
                # Kategorie/Fachbereich
                if 'categories' in volume_info and volume_info['categories']:
                    metadata['categories'] = volume_info['categories']
                
                # Sprache
                if 'language' in volume_info:
                    metadata['language'] = volume_info['language']
                
                return metadata
        
        return metadata
    
    def _fetch_openalex_metadata(self, title: str, authors: List[str]) -> Dict[str, Any]:
        """Ruft Metadaten von OpenAlex ab."""
        metadata = {}
        
        if not title:
            return metadata
        
        # Bereite die Suche vor
        query = f"title.search:{title}"
        if authors and authors[0]:
            # Versuche, Vor- und Nachnamen zu trennen
            name_parts = authors[0].split()
            if len(name_parts) > 1:
                last_name = name_parts[-1]
                query += f"&filter=author.display_name.search:{last_name}"
        
        url = f"https://api.openalex.org/works?filter={query}&sort=relevance_score:desc&per-page=1"
        headers = {"Accept": "application/json"}
        
        response = requests.get(url, headers=headers)
        
        if response.status_code == 200:
            data = response.json()
            if 'results' in data and data['results']:
                work = data['results'][0]
                
                # Titel
                if 'title' in work:
                    metadata['title'] = work['title']
                
                # DOI
                if 'doi' in work and work['doi']:
                    metadata['doi'] = work['doi']
                
                # Autoren
                if 'authorships' in work:
                    authors = []
                    for authorship in work['authorships']:
                        if 'author' in authorship and 'display_name' in authorship['author']:
                            authors.append(authorship['author']['display_name'])
                    if authors:
                        metadata['author'] = authors
                
                # Jahr
                if 'publication_year' in work:
                    metadata['year'] = work['publication_year']
                
                # Journal/Quelle
                if 'primary_location' in work and 'source' in work['primary_location']:
                    source = work['primary_location']['source']
                    if 'display_name' in source:
                        metadata['journal'] = source['display_name']
                
                # Zitationen
                if 'cited_by_count' in work:
                    metadata['cited_by_count'] = work['cited_by_count']
                
                # Fachgebiet
                if 'concepts' in work and work['concepts']:
                    concepts = []
                    for concept in work['concepts']:
                        if 'display_name' in concept:
                            concepts.append(concept['display_name'])
                    if concepts:
                        metadata['concepts'] = concepts
                
                # Open Access Status
                if 'open_access' in work and 'is_oa' in work['open_access']:
                    metadata['is_open_access'] = work['open_access']['is_oa']
                
                return metadata
        
        return metadata
    
    def split_text_into_chunks(self, text: str, language: str = "auto") -> List[Dict[str, Any]]:
        """
        Teilt den Text in Chunks auf für die Vektorisierung.
        
        Args:
            text: Der zu teilende Text
            language: Die Sprache des Textes ("de", "en", "auto")
            
        Returns:
            Liste von Chunks mit Text und Metadaten
        """
        # Wenn Sprache "auto", versuche die Sprache zu erkennen
        if language == "auto":
            language = self._detect_language(text)
        
        # NLP-Modell für die entsprechende Sprache laden
        nlp = self._load_spacy_model(language)
        
        # Chunks mit RecursiveCharacterTextSplitter erstellen
        raw_chunks = self.text_splitter.split_text(text)
        
        # Chunks mit Metadaten anreichern
        chunks = []
        for i, chunk_text in enumerate(raw_chunks):
            # Spacy-Verarbeitung für zusätzliche Metadaten (Entitäten, etc.)
            doc = nlp(chunk_text[:10000])  # Limit für Performance
            
            # Wichtige Entitäten extrahieren
            entities = {}
            for ent in doc.ents:
                entity_type = ent.label_
                if entity_type not in entities:
                    entities[entity_type] = []
                if ent.text not in entities[entity_type]:
                    entities[entity_type].append(ent.text)
            
            # Chunk mit Metadaten erstellen
            chunk = {
                "chunk_id": i,
                "text": chunk_text,
                "entities": entities,
                "language": language,
                "token_count": len(doc)
            }
            
            chunks.append(chunk)
        
        return chunks
    
    def _detect_language(self, text: str) -> str:
        """Erkennt die Sprache eines Textes (vereinfachte Version)."""
        # Eine einfache sprachunabhängige Erkennung basierend auf häufigen Wörtern
        text_sample = text[:1000].lower()
        
        # Zähle deutsche und englische häufige Wörter
        de_words = ["der", "die", "das", "und", "ist", "von", "für", "auf", "mit", "dem", "sich", "des", "ein", "nicht", "auch"]
        en_words = ["the", "and", "of", "to", "in", "is", "that", "for", "it", "as", "was", "with", "be", "by", "on"]
        
        de_count = sum(1 for word in de_words if f" {word} " in f" {text_sample} ")
        en_count = sum(1 for word in en_words if f" {word} " in f" {text_sample} ")
        
        return "de" if de_count > en_count else "en"
    
    def _save_processing_results(self, doc_dir: str, text: str, metadata: Dict[str, Any], chunks: List[Dict[str, Any]]):
        """Speichert die Verarbeitungsergebnisse."""
        # Text speichern
        with open(os.path.join(doc_dir, "fulltext.txt"), "w", encoding="utf-8") as f:
            f.write(text)
        
        # Metadaten speichern
        with open(os.path.join(doc_dir, "metadata.json"), "w", encoding="utf-8") as f:
            json.dump(metadata, f, ensure_ascii=False, indent=2)
        
        # Chunks speichern
        with open(os.path.join(doc_dir, "chunks.json"), "w", encoding="utf-8") as f:
            json.dump(chunks, f, ensure_ascii=False, indent=2)


# Test-Funktion zum einfachen Testen des Prozessors
def test_processor(filepath: str):
    """Testet den Dokumentenprozessor mit einer Datei."""
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    upload_dir = os.path.join(base_dir, "uploads")
    processed_dir = os.path.join(base_dir, "processed")
    
    processor = DocumentProcessor(upload_dir, processed_dir)
    
    # Kopiere die Datei ins Upload-Verzeichnis, falls sie nicht dort ist
    filename = os.path.basename(filepath)
    upload_path = os.path.join(upload_dir, filename)
    
    if filepath != upload_path:
        os.makedirs(upload_dir, exist_ok=True)
        shutil.copy2(filepath, upload_path)
    
    # Verarbeite die Datei
    result = processor.process_document(filename)
    
    print(f"Verarbeitung abgeschlossen: {result}")
    return result