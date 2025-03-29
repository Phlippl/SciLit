"""
PPTX-Parser für SciLit
-------------------
Parser für Microsoft PowerPoint-Präsentationen.
"""

import logging
from typing import Dict, Tuple, Any

# PPTX-Verarbeitung
try:
    import pptx
except ImportError:
    pptx = None

from app.core.document.parsers.base_parser import DocumentParser, DocumentParsingError

# Logger konfigurieren
logger = logging.getLogger("scilit.document.parsers.pptx")

class PPTXParser(DocumentParser):
    """
    Parser für PPTX-Dokumente (PowerPoint).
    
    Dieser Parser extrahiert Text und Metadaten aus PPTX-Dateien
    mit dem python-pptx Paket.
    """
    
    def parse(self, filepath: str) -> Tuple[str, Dict[str, Any]]:
        """
        Parst ein PPTX-Dokument und extrahiert Text und Metadaten.
        
        Args:
            filepath: Pfad zum PPTX
            
        Returns:
            Tuple aus (extrahierter Text, Dictionary mit Metadaten)
            
        Raises:
            DocumentParsingError: Bei Problemen mit dem Parsing
        """
        if not pptx:
            raise DocumentParsingError("python-pptx ist nicht installiert")
        
        try:
            # Metadaten aus Dateiinformationen
            metadata = self._extract_basic_metadata(filepath)
            
            # PPTX öffnen
            presentation = pptx.Presentation(filepath)
            
            # Text aus Folien extrahieren
            text_slides = []
            
            for i, slide in enumerate(presentation.slides):
                slide_text = []
                slide_text.append(f"--- Folie {i+1} ---")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                
                if slide_text:
                    text_slides.append("\n".join(slide_text))
            
            # Metadaten aus PPTX-Eigenschaften
            if hasattr(presentation.core_properties, "title") and presentation.core_properties.title:
                metadata["title"] = presentation.core_properties.title
                
            if hasattr(presentation.core_properties, "author") and presentation.core_properties.author:
                metadata["author"] = [author.strip() for author in presentation.core_properties.author.split(';')]
                
            if hasattr(presentation.core_properties, "created") and presentation.core_properties.created:
                metadata["created"] = presentation.core_properties.created.isoformat()
                metadata["year"] = presentation.core_properties.created.year
                
            # Seitenzahl = Anzahl der Folien
            self.page_count = len(presentation.slides)
            metadata["page_count"] = self.page_count
            
            return self._clean_text("\n\n".join(text_slides)), metadata
            
        except Exception as e:
            raise DocumentParsingError(f"Fehler beim PPTX-Parsing: {str(e)}")