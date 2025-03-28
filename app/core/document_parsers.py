"""
Document Parsers für SciLit
--------------------------
Enthält Parser für verschiedene Dateiformate.
"""

import os
import re
import logging
import tempfile
from typing import Dict, List, Optional, Tuple, Any, Protocol
from abc import ABC, abstractmethod

# Logger konfigurieren
logger = logging.getLogger("scilit.document_parsers")

class DocumentParsingError(Exception):
    """Fehler bei der Dokumentenverarbeitung."""
    pass

class DocumentParser(ABC):
    """Abstrakte Basisklasse für Dokumentenparser."""
    
    @abstractmethod
    def parse(self, filepath: str) -> Tuple[str, Dict[str, Any]]:
        """
        Extrahiert Text und Metadaten aus einem Dokument.
        
        Args:
            filepath: Pfad zur Datei
            
        Returns:
            Tuple aus (extrahierter Text, Dictionary mit Metadaten)
        """
        pass

class PDFParser(DocumentParser):
    """Parser für PDF-Dateien."""
    
    def __init__(self, ocr_if_needed: bool = True, ocr_language: str = "auto"):
        self.ocr_if_needed = ocr_if_needed
        self.ocr_language = ocr_language
        
        # OCR-Sprachkonfiguration
        self.OCR_LANGUAGES = {
            'de': 'deu',
            'en': 'eng',
            'auto': 'deu+eng'  # Mehrsprachenerkennung
        }
    
    def parse(self, filepath: str) -> Tuple[str, Dict[str, Any]]:
        """Verarbeitet eine PDF-Datei."""
        import PyPDF2
        
        text = ""
        metadata = {}
        
        try:
            with open(filepath, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                # Metadaten aus PDF extrahieren
                pdf_info = pdf_reader.metadata
                if pdf_info:
                    # Standardmetadaten extrahieren
                    for key in pdf_info:
                        clean_key = key.strip('/').lower()
                        metadata[clean_key] = pdf_info[key]
                
                # Text aus allen Seiten extrahieren
                page_texts = []
                needs_ocr = False
                
                for i, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    
                    # Überprüfen, ob OCR benötigt wird (wenn Seite nicht genug extrahierbaren Text hat)
                    if not page_text or len(page_text.strip()) < 100:
                        if self.ocr_if_needed:
                            logger.info(f"Seite {i+1} hat wenig Text - versuche OCR")
                            needs_ocr = True
                            break
                    
                    page_texts.append(page_text)
                
                # OCR anwenden, wenn nötig und aktiviert
                if needs_ocr and self.ocr_if_needed:
                    logger.info(f"Wende OCR auf {filepath} an")
                    text = self._apply_ocr_to_pdf(filepath)
                else:
                    text = "\n\n".join(page_texts)
                
                # Versuche, das Publikationsjahr zu extrahieren
                if 'creationdate' in metadata:
                    year_match = re.search(r'D:(\d{4})', metadata['creationdate'])
                    if year_match:
                        metadata['year'] = int(year_match.group(1))
                
                # Seitenzahl hinzufügen
                metadata['page_count'] = len(pdf_reader.pages)
                
        except Exception as e:
            logger.error(f"Fehler beim Verarbeiten von PDF {filepath}: {str(e)}")
            raise DocumentParsingError(f"Fehler beim Verarbeiten von PDF: {str(e)}")
        
        return text, metadata
    
    def _apply_ocr_to_pdf(self, filepath: str) -> str:
        """Wendet OCR auf eine PDF-Datei an."""
        from PIL import Image
        import pytesseract
        
        logger.info(f"Starte OCR-Verarbeitung für {filepath}")
        
        # Versuche zuerst mit PyMuPDF (fitz)
        try:
            import fitz  # PyMuPDF
            
            ocr_lang = self.OCR_LANGUAGES.get(self.ocr_language, self.OCR_LANGUAGES['auto'])
            doc = fitz.open(filepath)
            text_pages = []
            
            with tempfile.TemporaryDirectory() as temp_dir:
                for page_num, page in enumerate(doc):
                    logger.info(f"Verarbeite Seite {page_num+1}/{len(doc)}")
                    # Bild der Seite extrahieren
                    pix = page.get_pixmap()
                    img_path = f"{temp_dir}/page_{page_num+1}.png"
                    pix.save(img_path)
                    
                    # OCR anwenden
                    try:
                        img = Image.open(img_path)
                        page_text = pytesseract.image_to_string(img, lang=ocr_lang)
                        text_pages.append(page_text)
                    except Exception as e:
                        logger.error(f"OCR-Fehler bei Seite {page_num+1}: {str(e)}")
                        text_pages.append(f"[OCR-FEHLER: {str(e)}]")
            
            return "\n\n".join(text_pages)
            
        except ImportError as e:
            logger.warning(f"PyMuPDF (fitz) nicht verfügbar: {str(e)}. Versuche alternative Methode...")
            
        # Alternative Methode mit pdf2image
        try:
            from pdf2image import convert_from_path
            
            ocr_lang = self.OCR_LANGUAGES.get(self.ocr_language, self.OCR_LANGUAGES['auto'])
            text_pages = []
            
            with tempfile.TemporaryDirectory() as temp_dir:
                try:
                    # PDF in Bilder umwandeln
                    logger.info(f"Konvertiere PDF zu Bildern mit pdf2image...")
                    pages = convert_from_path(filepath, dpi=300, output_folder=temp_dir)
                    
                    for i, page in enumerate(pages):
                        logger.info(f"OCR für Seite {i+1}/{len(pages)}")
                        # OCR für jede Seite anwenden
                        try:
                            page_text = pytesseract.image_to_string(page, lang=ocr_lang)
                            text_pages.append(page_text)
                        except Exception as e:
                            logger.error(f"OCR-Fehler bei Seite {i+1}: {str(e)}")
                            text_pages.append(f"[OCR-FEHLER: {str(e)}]")
                    
                    return "\n\n".join(text_pages)
                except Exception as e:
                    logger.error(f"Fehler bei pdf2image: {str(e)}")
                    return f"[FEHLER BEI DER OCR-VERARBEITUNG: {str(e)}. Bitte stelle sicher, dass Poppler korrekt installiert ist.]"
                    
        except ImportError as e:
            error_msg = f"Weder PyMuPDF (fitz) noch pdf2image sind verfügbar: {str(e)}"
            logger.error(error_msg)
            return f"[OCR NICHT VERFÜGBAR: {error_msg}. Bitte installiere PyMuPDF oder pdf2image.]"


class DOCXParser(DocumentParser):
    """Parser für DOCX-Dateien."""
    
    def parse(self, filepath: str) -> Tuple[str, Dict[str, Any]]:
        """Verarbeitet eine DOCX-Datei."""
        import docx
        
        doc = docx.Document(filepath)
        text = "\n".join([para.text for para in doc.paragraphs])
        
        # Metadaten extrahieren
        metadata = {}
        core_properties = doc.core_properties
        
        if core_properties.title:
            metadata['title'] = core_properties.title
        if core_properties.author:
            metadata['author'] = core_properties.author
        if core_properties.created:
            metadata['created'] = core_properties.created.isoformat()
            metadata['year'] = core_properties.created.year
        
        return text, metadata


class TXTParser(DocumentParser):
    """Parser für TXT-Dateien."""
    
    def parse(self, filepath: str) -> Tuple[str, Dict[str, Any]]:
        """Verarbeitet eine TXT-Datei."""
        with open(filepath, 'r', encoding='utf-8', errors='replace') as file:
            text = file.read()
        
        # Bei TXT-Dateien müssen wir Metadaten aus dem Text extrahieren
        metadata = {}
        
        # Erste nicht-leere Zeile könnte der Titel sein
        lines = [line.strip() for line in text.split('\n') if line.strip()]
        if lines:
            metadata['title'] = lines[0]
        
        return text, metadata


class EPUBParser(DocumentParser):
    """Parser für EPUB-Dateien."""
    
    def parse(self, filepath: str) -> Tuple[str, Dict[str, Any]]:
        """Verarbeitet eine EPUB-Datei."""
        import epub2txt
        
        # epub2txt extrahiert Text aus EPUB-Dateien
        text = epub2txt.epub2txt(filepath)
        
        # Metadaten-Extraktion aus EPUB erfordert zusätzliche Bibliotheken
        metadata = {}
        
        try:
            import ebooklib
            from ebooklib import epub
            
            book = epub.read_epub(filepath)
            
            # Metadaten extrahieren
            if book.metadata:
                # Dublin Core Metadaten
                dc_metadata = book.metadata.get('http://purl.org/dc/elements/1.1/', {})
                
                # Titel
                if 'title' in dc_metadata:
                    metadata['title'] = dc_metadata['title'][0][0]
                
                # Autoren
                if 'creator' in dc_metadata:
                    authors = [creator[0] for creator in dc_metadata['creator']]
                    metadata['author'] = authors
                
                # Publikationsdatum
                if 'date' in dc_metadata:
                    date_str = dc_metadata['date'][0][0]
                    # Versuche, Jahr zu extrahieren
                    year_match = re.search(r'(\d{4})', date_str)
                    if year_match:
                        metadata['year'] = int(year_match.group(1))
                
                # Publisher
                if 'publisher' in dc_metadata:
                    metadata['publisher'] = dc_metadata['publisher'][0][0]
        
        except ImportError:
            logger.warning("ebooklib nicht installiert, beschränkte Metadaten-Extraktion")
        
        return text, metadata


class PPTXParser(DocumentParser):
    """Parser für PPTX-Dateien."""
    
    def parse(self, filepath: str) -> Tuple[str, Dict[str, Any]]:
        """Verarbeitet eine PPTX-Datei."""
        try:
            from pptx import Presentation
            
            presentation = Presentation(filepath)
            
            # Text aus allen Folien extrahieren
            text_runs = []
            
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            
            text = "\n\n".join(text_runs)
            
            # Metadaten extrahieren
            metadata = {}
            core_properties = presentation.core_properties
            
            if core_properties.title:
                metadata['title'] = core_properties.title
            if core_properties.author:
                metadata['author'] = core_properties.author
            if core_properties.created:
                metadata['created'] = core_properties.created.isoformat()
                metadata['year'] = core_properties.created.year
            
            # Anzahl der Folien
            metadata['slide_count'] = len(presentation.slides)
            
            return text, metadata
            
        except ImportError:
            logger.error("python-pptx nicht installiert")
            raise ImportError("Zur Verarbeitung von PPTX-Dateien wird python-pptx benötigt")


def determine_parser_for_file(filepath: str, ocr_if_needed: bool = True) -> Optional[DocumentParser]:
    """
    Bestimmt den geeigneten Parser für eine Datei basierend auf dem Dateityp.
    
    Args:
        filepath: Pfad zur Datei
        ocr_if_needed: Ob OCR bei PDF-Dateien bei Bedarf angewendet werden soll
        
    Returns:
        Ein DocumentParser-Objekt oder None, wenn kein passender Parser gefunden wird
    """
    file_ext = os.path.splitext(filepath)[1].lower()
    
    if file_ext == '.pdf':
        return PDFParser(ocr_if_needed=ocr_if_needed)
    elif file_ext == '.docx':
        return DOCXParser()
    elif file_ext == '.txt':
        return TXTParser()
    elif file_ext == '.epub':
        return EPUBParser()
    elif file_ext in ['.pptx', '.ppt']:
        return PPTXParser()
    else:
        logger.warning(f"Dateityp {file_ext} wird nicht unterstützt")
        return None